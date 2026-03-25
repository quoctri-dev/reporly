"""
Reporly — PPTX Report Exporter
Generates professional PowerPoint report from Report dataclass.
"""
import io
import logging
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from src.templates import get_template

logger = logging.getLogger(__name__)


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def export_pptx(report, template_name: str = "minimal", app_name: str = "Reporly") -> bytes:
    """
    Export Report to PPTX bytes.

    Args:
        report: Report dataclass (from core.models)
        template_name: Template style name
        app_name: App name for branding

    Returns:
        PPTX file as bytes
    """
    tmpl = get_template(template_name)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Cover ---
    _add_cover_slide(prs, report, tmpl, app_name)

    # --- Slide 2: Data Overview ---
    _add_overview_slide(prs, report.data_profile, tmpl)

    # --- Slide 3: Key Insights ---
    if report.insights:
        _add_insights_slide(prs, report.insights, tmpl)

    # --- Slide 4-N: Charts ---
    for i, chart_bytes in enumerate(report.charts):
        _add_chart_slide(prs, chart_bytes, tmpl, i + 1)

    # --- Last slide: Footer ---
    _add_footer_slide(prs, report, tmpl, app_name)

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    logger.info("PPTX exported: %.1f KB", len(pptx_bytes) / 1024)
    return pptx_bytes


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _add_cover_slide(prs, report, tmpl, app_name: str):
    """Title slide with report name, filename, date."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background accent bar (top)
    _add_rect(slide, 0, 0, prs.slide_width, Inches(0.15), tmpl.primary_color)

    # Title
    _add_textbox(
        slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
        report.title,
        font_name=tmpl.font_heading, font_size=Pt(40), bold=True,
        color=tmpl.primary_color, alignment=PP_ALIGN.LEFT,
    )

    # Subtitle: filename + dimensions
    subtitle = (
        f"{report.data_profile.filename}  |  "
        f"{report.data_profile.rows:,} rows x {report.data_profile.columns} columns"
    )
    _add_textbox(
        slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
        subtitle,
        font_name=tmpl.font_body, font_size=Pt(18),
        color=tmpl.secondary_color, alignment=PP_ALIGN.LEFT,
    )

    # Date
    _add_textbox(
        slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
        report.generated_at.strftime("%Y-%m-%d %H:%M"),
        font_name=tmpl.font_body, font_size=Pt(14),
        color=tmpl.secondary_color, alignment=PP_ALIGN.LEFT,
    )

    # Bottom bar
    _add_rect(slide, 0, prs.slide_height - Inches(0.15), prs.slide_width, Inches(0.15), tmpl.accent_color)


def _add_overview_slide(prs, profile, tmpl):
    """Data overview table slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, prs.slide_width, Inches(0.08), tmpl.primary_color)

    _add_textbox(
        slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
        "Data Overview",
        font_name=tmpl.font_heading, font_size=Pt(28), bold=True,
        color=tmpl.primary_color,
    )

    # Summary stats table
    rows_data = [
        ["Metric", "Value"],
        ["Total Rows", f"{profile.rows:,}"],
        ["Total Columns", str(profile.columns)],
        ["Numeric Columns", str(profile.basic_stats.get("numeric_columns", 0))],
        ["Categorical Columns", str(profile.basic_stats.get("categorical_columns", 0))],
        ["Missing Values", f"{profile.basic_stats.get('total_nulls', 0):,} ({profile.basic_stats.get('null_pct', 0)}%)"],
        ["Memory Size", f"{profile.basic_stats.get('memory_mb', 0)} MB"],
    ]

    n_rows = len(rows_data)
    n_cols = 2
    tbl = slide.shapes.add_table(n_rows, n_cols, Inches(0.8), Inches(1.3), Inches(6), Inches(0.4 * n_rows)).table
    tbl.columns[0].width = Inches(3)
    tbl.columns[1].width = Inches(3)

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.name = tmpl.font_body
            if r == 0:
                para.font.bold = True
                para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _hex_to_rgb(tmpl.primary_color)
            else:
                para.font.color.rgb = _hex_to_rgb("#374151")

    # Column details (compact, right side)
    col_info = profile.column_info[:12]  # Max 12 in slide
    if col_info:
        ci_rows = [["Column", "Type", "Nulls"]]
        for ci in col_info:
            ci_rows.append([ci.name[:20], ci.dtype, f"{ci.null_pct}%"])

        ci_n = len(ci_rows)
        tbl2 = slide.shapes.add_table(ci_n, 3, Inches(7.5), Inches(1.3), Inches(5), Inches(0.35 * ci_n)).table
        tbl2.columns[0].width = Inches(2.2)
        tbl2.columns[1].width = Inches(1.4)
        tbl2.columns[2].width = Inches(1.4)

        for r, row in enumerate(ci_rows):
            for c, val in enumerate(row):
                cell = tbl2.cell(r, c)
                cell.text = val
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10)
                para.font.name = tmpl.font_body
                if r == 0:
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _hex_to_rgb(tmpl.secondary_color)
                else:
                    para.font.color.rgb = _hex_to_rgb("#374151")


def _add_insights_slide(prs, insights, tmpl):
    """Key insights as bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, prs.slide_width, Inches(0.08), tmpl.primary_color)

    _add_textbox(
        slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
        "Key Insights",
        font_name=tmpl.font_heading, font_size=Pt(28), bold=True,
        color=tmpl.primary_color,
    )

    importance_colors = {"high": "#DC2626", "medium": "#D97706", "low": "#059669"}

    y = Inches(1.4)
    for i, insight in enumerate(insights):
        color = importance_colors.get(insight.importance, "#6B7280")

        # Importance badge
        _add_textbox(
            slide, Inches(0.8), y, Inches(1), Inches(0.35),
            f"[{insight.importance.upper()}]",
            font_name=tmpl.font_body, font_size=Pt(10), bold=True,
            color=color,
        )

        # Title
        _add_textbox(
            slide, Inches(2), y, Inches(10), Inches(0.35),
            f"{i + 1}. {insight.title}",
            font_name=tmpl.font_heading, font_size=Pt(16), bold=True,
            color="#1F2937",
        )

        # Description
        _add_textbox(
            slide, Inches(2), y + Inches(0.4), Inches(10), Inches(0.6),
            insight.description,
            font_name=tmpl.font_body, font_size=Pt(12),
            color="#6B7280",
        )

        y += Inches(1.1)
        if y > Inches(6.5):
            break


def _add_chart_slide(prs, chart_bytes: bytes, tmpl, chart_num: int):
    """One chart per slide, full width."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, prs.slide_width, Inches(0.08), tmpl.primary_color)

    _add_textbox(
        slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
        f"Chart {chart_num}",
        font_name=tmpl.font_heading, font_size=Pt(20), bold=True,
        color=tmpl.primary_color,
    )

    img_stream = io.BytesIO(chart_bytes)
    slide.shapes.add_picture(img_stream, Inches(1.5), Inches(1.2), Inches(10), Inches(5.5))


def _add_footer_slide(prs, report, tmpl, app_name: str):
    """Final slide with branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, tmpl.primary_color)

    _add_textbox(
        slide, Inches(2), Inches(2.5), Inches(9), Inches(1),
        f"Generated by {app_name}",
        font_name=tmpl.font_heading, font_size=Pt(32), bold=True,
        color="#FFFFFF", alignment=PP_ALIGN.CENTER,
    )

    _add_textbox(
        slide, Inches(2), Inches(3.8), Inches(9), Inches(0.6),
        report.generated_at.strftime("%Y-%m-%d %H:%M"),
        font_name=tmpl.font_body, font_size=Pt(16),
        color="#FFFFFF", alignment=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _add_textbox(slide, left, top, width, height, text, *,
                 font_name="Helvetica", font_size=Pt(12), bold=False,
                 color="#000000", alignment=PP_ALIGN.LEFT):
    """Add a text box to slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = _hex_to_rgb(color) if isinstance(color, str) else color
    p.alignment = alignment


def _add_rect(slide, left, top, width, height, fill_color: str):
    """Add a filled rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    shape.line.fill.background()
